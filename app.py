import streamlit as st
import base64
import openai
import pptx
from pptx.util import Inches, Pt
import os

from dotenv import load_dotenv
load_dotenv()

openai.api_key = os.getenv('OPENAI_API_KEY')  # Replace with your actual API key

# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)


def generate_slide_titles(topic):
    prompt = f"Generate 5 slide titles for the topic '{topic}'."
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=200,
    )
    return response['choices'][0]['text'].split("\n")

def generate_slide_content(slide_title):
    prompt = f"Generate content for the slide: '{slide_title}'."
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=500,  # Adjust as needed based on the desired content length
    )
    return response['choices'][0]['text']


def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        # Customize font size for titles and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    prs.save(f"generated_ppt/{topic}_presentation.pptx")
    

def main():
    st.title("PowerPoint Presentation Generator with GPT-3.5-turbo")

    topic = st.text_input("Enter the topic for your presentation:")
    generate_button = st.button("Generate Presentation")

    if generate_button and topic:
        st.info("Generating presentation... Please wait.")
        slide_titles = generate_slide_titles(topic)
        filtered_slide_titles= [item for item in slide_titles if item.strip() != '']
        print("Slide Title: ", filtered_slide_titles)
        slide_contents = [generate_slide_content(title) for title in filtered_slide_titles]
        print("Slide Contents: ", slide_contents)
        create_presentation(topic, filtered_slide_titles, slide_contents)
        print("Presentation generated successfully!")

        st.success("Presentation generated successfully!")
        st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)

def get_ppt_download_link(topic):
    ppt_filename = f"generated_ppt/{topic}_presentation.pptx"

    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'


if __name__ == "__main__":
    main()
