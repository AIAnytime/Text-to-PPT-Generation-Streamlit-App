import openai
import pptx
from pptx.util import Inches
import os
import time

from dotenv import load_dotenv
load_dotenv()

openai.api_key = os.getenv('OPENAI_API_KEY')  # Replace with your actual API key

def generate_slide_titles(topic):
    prompt = f"Generate 10 slide titles for the topic '{topic}'."
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=200,
    )
    return response['choices'][0]['text'].split("\n")

def generate_slide_content(slide_title):
    prompt = f"Generate content for the slide: '{slide_title}'."
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=500,  # Adjust as needed based on the desired content length
    )
    return response['choices'][0]['text']


def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

    prs.save(f"{topic}_presentation.pptx")

def main():
    start = time.time()
    topic = "Explainable AI"  # Replace with your desired topic

    slide_titles = generate_slide_titles(topic)
    filtered_slide_titles= [item for item in slide_titles if item.strip() != '']
    print("Slide Title: ", filtered_slide_titles)
    slide_contents = [generate_slide_content(title) for title in filtered_slide_titles]
    print("Slide Contents: ", slide_contents)
    create_presentation(topic, filtered_slide_titles, slide_contents)
    end = time.time()

    print("Total Duration: ", round(end - start, 2), " secs")
    print("Presentation generated successfully!")
    
if __name__ == "__main__":
    main()
